import requests
import os
import shutil
from PIL import Image
from io import BytesIO
import time
from tqdm import tqdm
import cairosvg
import argparse
from pptx import Presentation
from pptx.util import Inches

def download_and_convert_slides(base_url, output_dir, start_slide=1, max_slides=100, format="jpeg"):
    """
    Download SVG slides and convert them to the specified format.
    
    Args:
        base_url (str): Base URL for slides (without slide number and extension)
        output_dir (str): Directory to save the converted images
        start_slide (int): Starting slide number
        max_slides (int): Maximum number of slides to attempt downloading
        format (str): Output format (jpeg, png)
    
    Returns:
        list: Paths to successfully downloaded and converted images
    """
    # Create output directory if it doesn't exist
    os.makedirs(output_dir, exist_ok=True)
    
    # Prepare URL (ensure it ends with 'slide' without number or extension)
    if not base_url.endswith('slide'):
        if base_url.endswith('/'):
            base_url = base_url + 'slide'
        elif base_url.endswith('.svg'):
            # Extract base URL from a complete URL like slide1.svg
            base_url = base_url.rsplit('slide', 1)[0] + 'slide'
        else:
            base_url = base_url + '/slide'
    
    downloaded_files = []
    consecutive_failures = 0
    max_consecutive_failures = 3
    
    print(f"Starting download from: {base_url}1.svg")
    
    # Use tqdm for a progress bar
    progress_bar = tqdm(range(start_slide, start_slide + max_slides))
    
    for slide_num in progress_bar:
        if consecutive_failures >= max_consecutive_failures:
            print(f"\nStopping after {max_consecutive_failures} consecutive failures")
            break
        
        slide_url = f"{base_url}{slide_num}.svg"
        output_file = os.path.join(output_dir, f"slide{slide_num}.{format}")
        
        progress_bar.set_description(f"Downloading slide {slide_num}")
        
        try:
            # Download SVG
            response = requests.get(slide_url, timeout=10)
            response.raise_for_status()  # Raise exception for 404s or other errors
            
            # Convert SVG to desired format
            if format == "jpeg" or format == "jpg":
                # Use cairosvg to convert SVG to PNG
                png_data = cairosvg.svg2png(bytestring=response.content)
                # Use PIL to convert PNG to JPEG
                image = Image.open(BytesIO(png_data))
                rgb_image = image.convert('RGB')  # Convert to RGB for JPEG
                rgb_image.save(output_file, quality=90)
            elif format == "png":
                png_data = cairosvg.svg2png(bytestring=response.content)
                with open(output_file, 'wb') as f:
                    f.write(png_data)
            else:
                raise ValueError(f"Unsupported format: {format}")
            
            downloaded_files.append(output_file)
            consecutive_failures = 0
            
        except requests.exceptions.HTTPError as e:
            if e.response.status_code == 404:
                consecutive_failures += 1
                progress_bar.set_description(f"Slide {slide_num} not found")
            else:
                print(f"\nError downloading {slide_url}: {e}")
                consecutive_failures += 1
        except Exception as e:
            print(f"\nError processing {slide_url}: {e}")
            consecutive_failures += 1
        
        # Small delay to avoid overwhelming the server
        time.sleep(0.5)
    
    print(f"\nSuccessfully downloaded {len(downloaded_files)} slides")
    return downloaded_files

def create_presentation(image_files, output_file):
    """
    Create a PowerPoint presentation with the downloaded images.
    
    Args:
        image_files (list): List of image file paths
        output_file (str): Path to save the PowerPoint presentation
    """
    prs = Presentation()
    
    # Set slide dimensions (16:9 aspect ratio)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    for image_file in tqdm(image_files, desc="Creating presentation"):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide
        
        # Calculate image dimensions to fit the slide
        img = Image.open(image_file)
        width, height = img.size
        
        # Calculate aspect ratio
        img_ratio = width / height
        slide_ratio = prs.slide_width / prs.slide_height
        
        if img_ratio > slide_ratio:
            # Image is wider than slide
            pic_width = prs.slide_width
            pic_height = pic_width / img_ratio
            left = 0
            top = (prs.slide_height - pic_height) / 2
        else:
            # Image is taller than slide
            pic_height = prs.slide_height
            pic_width = pic_height * img_ratio
            top = 0
            left = (prs.slide_width - pic_width) / 2
        
        # Add picture to slide
        slide.shapes.add_picture(image_file, left, top, width=pic_width, height=pic_height)
    
    prs.save(output_file)
    print(f"Presentation saved to {output_file}")

def cleanup_temp_files(output_dir):
    """
    Delete the temporary directory with downloaded slide images.
    
    Args:
        output_dir (str): Directory to remove
    """
    try:
        if os.path.exists(output_dir):
            shutil.rmtree(output_dir)
            print(f"Cleaned up temporary files: {output_dir} directory deleted")
    except Exception as e:
        print(f"Warning: Could not delete temporary directory {output_dir}: {e}")

def main():
    parser = argparse.ArgumentParser(description='Download slides and create a presentation')
    parser.add_argument('base_url', help='Base URL for slides (e.g., https://example.com/presentation/id/svgs/slide)')
    parser.add_argument('--output-dir', default='slides', help='Directory to save the images (default: slides)')
    parser.add_argument('--start', type=int, default=1, help='Starting slide number (default: 1)')
    parser.add_argument('--max', type=int, default=100, help='Maximum number of slides to download (default: 100)')
    parser.add_argument('--format', choices=['jpeg', 'png'], default='jpeg', help='Output image format (default: jpeg)')
    parser.add_argument('--presentation', default='presentation.pptx', help='Output PowerPoint file (default: presentation.pptx)')
    parser.add_argument('--keep-images', action='store_true', help='Keep the image files after creating the presentation')
    
    args = parser.parse_args()
    
    print("Slide Downloader and Presentation Creator")
    print("=========================================")
    
    # Download and convert slides
    image_files = download_and_convert_slides(
        args.base_url, 
        args.output_dir,
        args.start,
        args.max,
        args.format
    )
    
    if image_files:
        # Create presentation
        create_presentation(image_files, args.presentation)
        
        # Clean up temporary files unless --keep-images is specified
        if not args.keep_images:
            cleanup_temp_files(args.output_dir)

if __name__ == "__main__":
    main()